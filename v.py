from pptx import Presentation
from pptx.util import Inches, Pt

# Crear presentación
prs = Presentation()

# --- Diapositiva 1: Título ---
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
slide_1.shapes.title.text = "Análisis de Ventas - Supermercado"

# --- Diapositiva 2: Descripción general ---
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
slide_2.shapes.title.text = "Descripción general"
slide_2.placeholders[1].text = (
    "El conjunto de datos recoge las ventas del supermercado entre el 1 de enero y el 4 de febrero de 2024.\n"
    "Cada registro representa una transacción individual con información sobre fechas, clientes y montos.\n"
    "En general, los datos son claros y no muestran errores visibles."
)

# --- Diapositiva 3: Atributos del conjunto ---
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
slide_3.shapes.title.text = "Atributos del conjunto"
slide_3.placeholders[1].text = (
    "- Número de pedido\n"
    "- Fechas de orden y envío\n"
    "- Nombre del cliente\n"
    "- Precio unitario\n"
    "- Cantidad comprada\n"
    "- Impuesto aplicado (10%)\n"
    "- Total final de la venta"
)

# --- Diapositiva 4: Patrones encontrados ---
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
slide_4.shapes.title.text = "Patrones encontrados"
slide_4.placeholders[1].text = (
    "• Los precios siguen valores fijos (19.99, 29.99, 49.99, etc.), lo que sugiere un catálogo definido.\n"
    "• Los productos de menor precio se compran en mayores cantidades.\n"
    "• Los artículos más costosos se adquieren en menor número.\n"
    "• El tiempo entre la compra y el envío suele ser de 1 a 3 días.\n"
    "• Algunos clientes repiten pedidos, reflejando compras frecuentes.\n"
    "• El volumen de ventas se mantiene estable sin picos notorios."
)

# --- Diapositiva 5: Conclusión ---
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
slide_5.shapes.title.text = "Conclusión"
slide_5.placeholders[1].text = (
    "El análisis muestra un sistema de ventas ordenado y constante.\n"
    "Los datos reflejan eficiencia en los envíos, clientes recurrentes "
    "y una relación clara entre los precios del catálogo y las cantidades compradas."
)

# Guardar la presentación
prs.save("Analisis_Supermercado_v2.pptx")
print("Presentación creada: Analisis_Supermercado_v2.pptx")
